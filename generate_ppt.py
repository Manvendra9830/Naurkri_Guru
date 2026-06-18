from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation(output_path, images_dir):
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Naukri_Guru\nAI-Powered Job Automation Platform"
    subtitle.text = "Manvendra Pratap Singh (CS22B1054)\nGuide: Dr. Neha Agarwal\nIIIT Raichur"

    # Slide 2: Introduction
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction & Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Job hunting is tedious, repetitive, and time-consuming."
    p = tf.add_paragraph()
    p.text = "High Application Volume: Hundreds of applications required to secure an offer."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Redundant Question Answering: Repeatedly filling similar screening questions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fragmented Outcome Tracking: Manual tracking of recruiter responses across platforms."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Naukri_Guru solves this by automating discovery, application, and tracking."
    p.level = 0

    # Slide 3: Objectives
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Project Objectives"
    tf = body_shape.text_frame
    tf.text = "Architect a stealth browser automation engine for LinkedIn."
    p = tf.add_paragraph()
    p.text = "Develop a confidence scoring algorithm for job suitability."
    p = tf.add_paragraph()
    p.text = "Build a multi-tier Q&A system (Config -> Memory -> AI with safety guards)."
    p = tf.add_paragraph()
    p.text = "Create an Indeed job scraper for multi-portal discovery."
    p = tf.add_paragraph()
    p.text = "Implement SQLite persistence and automated Gmail IMAP lifecycle tracking."
    p = tf.add_paragraph()
    p.text = "Engineer a cold email outreach pipeline for proactive recruiter engagement."

    # Helper function to add image slides
    def add_image_slide(title_text, img_filename):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = title_text
        img_path = os.path.join(images_dir, img_filename)
        if os.path.exists(img_path):
            # Calculate image position to center it
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11.333))
        else:
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = f"Image not found: {img_filename}"
        return slide

    # Slide 4: System Block Diagram
    add_image_slide("System Block Diagram", "block_diagram.png")

    # Slide 5: Software Architecture
    add_image_slide("Software Architecture", "architecture.png")

    # Slide 6: Module 1 - Stealth Browser Engine
    slide = add_image_slide("Module 1: Stealth Browser Automation", "module1_flow.png")
    
    # Slide 7: Module 2 - Job Discovery & Filtering
    slide = add_image_slide("Module 2: Job Discovery & Intelligent Filtering", "module2_flow.png")

    # Slide 8: Filtering Details
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Intelligent Filtering Funnel"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Multi-stage pipeline to narrow down candidate pool:"
    p = tf.add_paragraph()
    p.text = "Blacklist & Deduplication: Ignores blacklisted companies and already applied jobs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Confidence Scoring: Evaluates JD text (+intern, +fresher, -senior, -clearance)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Seniority & Degree Check: Bypasses PhD/Principal/Director requirements."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Experience Threshold: Matches candidate's experience against job requirements."
    p.level = 1

    # Slide 9: Module 3 - Intelligent Q&A
    slide = add_image_slide("Module 3: Intelligent Question Answering", "module3_flow.png")

    # Slide 10: Three-Tier Cascade
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Q&A Three-Tier Cascade Strategy"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Tier 1: Static Configuration"
    p = tf.add_paragraph()
    p.text = "Direct mapping for personal info (phone, address, etc.)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tier 2: Persistent Memory (memory.json)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Caches previously answered questions for future reuse. Skill-specific routing."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tier 3: AI Provider (Gemini, OpenAI, DeepSeek)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Generates answers for novel questions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Safety Guard function overrides hallucinated sensitive personal data."
    p.level = 1

    # Slide 11: Module 4 - Lifecycle & Persistence
    slide = add_image_slide("Module 4: Application Lifecycle & Persistence", "module4_flow.png")

    # Slide 12: Data Management & Export
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Data Management & Indeed Scraper"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Post-Submission Verification"
    p = tf.add_paragraph()
    p.text = "Verifies DOM for success markers before persisting to SQLite."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Gmail IMAP Synchronization"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Detects recruiter responses, matches with applications, updates status."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Indeed Scraper Module"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Discovers and archives Indeed jobs in a unified store with bad-word filtering."
    p.level = 1

    # Slide 13: Major Challenges
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Challenges & Solutions"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Challenge: LinkedIn DOM instability and A/B tests."
    p = tf.add_paragraph()
    p.text = "Solution: Multi-selector cascading strategy with robust fallback logic."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Challenge: AI hallucinating personal information."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Solution: Dedicated safety guard pattern matching."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Challenge: Duplicate job variants at the same company."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Solution: Session-level company+title deduplication cache."
    p.level = 1

    # Slide 14: Results Output
    add_image_slide("Results: Automation Output", "result1.png")

    # Slide 15: Exported Tracking
    add_image_slide("Results: Application Tracking (XLSX)", "result4.png")

    # Slide 16: Performance Metrics
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Performance Metrics & Evaluation"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Time Savings: ~90% reduction (15-20s per app vs. 3-5 mins manual)."
    p = tf.add_paragraph()
    p.text = "Success Rate: 87.9% application success (post-verification)."
    p = tf.add_paragraph()
    p.text = "Memory Convergence: Reaches autonomous operation within 3-5 sessions."
    p = tf.add_paragraph()
    p.text = "Filtering Efficacy: Removes ~79% of unsuitable listings automatically."
    p = tf.add_paragraph()
    p.text = "Email Classification: 85% match rate for recruiter responses via Gmail IMAP."
    
    # Slide 17: Demo Video (Placeholder)
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "[ PLACEHOLDER FOR DEMO VIDEO ]"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Slide 18: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Conclusion & Future Scope"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Conclusion"
    p = tf.add_paragraph()
    p.text = "Successfully automated the job application lifecycle."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Demonstrated stability with intelligent Q&A and multi-portal integration."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Scope"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Multi-Platform auto-apply (Naukri, Glassdoor, Wellfound)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Machine Learning-based confidence scoring from historical accept/reject data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Interactive analytics dashboard and interview scheduling integration."
    p.level = 1

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    images_dir = r"D:\Major_Project\Auto_job_applier_linkedIn\Naukri_Guru\Thesis\IIIT_Raichur_Thesis_Template_new\images"
    output_path = r"D:\Major_Project\Auto_job_applier_linkedIn\Naukri_Guru\Naukri_Guru_Thesis_Presentation.pptx"
    create_presentation(output_path, images_dir)
